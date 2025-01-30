import unittest
import pptx

from data_ingestion.presentation import Presentation, PresentationSlide

class PresentationTest(unittest.TestCase):
    def test_it_should_build_an_empty_presentation(self):
        presentation = Presentation()

        self.assertEquals(len(presentation.slides),0)


    def test_adding_a_first_slide(self):
        first_slide = PresentationSlide("1", "2001")

        presentation = Presentation()

        presentation.append_slide(first_slide)

        self.assertEquals(len(presentation.slides), 1)

    def test_import_of_pptx_with_text_from_ms_pptx(self):
        test_file_path = "src/test/test_data/offset-for-pages-test.pptx"
        pptx_presentation = pptx.Presentation(test_file_path)

        coscrad_presentation = Presentation.from_pptx(pptx_presentation)

        self.assertEquals(len(coscrad_presentation.slides), 4)

        first_slide = coscrad_presentation.slides[0]

        self.assertEquals("I have an offset for page numbers\nTrust me, I do",str(first_slide))

    def test_import_of_pptx_with_image(self):
        test_file_path = "src/test/test_data/test_presentation_with_images_in_seperate_shapes.pptx"
        pptx_presentation = pptx.Presentation(test_file_path)

        coscrad_presentation = Presentation.from_pptx(pptx_presentation)

        slide_with_image = coscrad_presentation.slides[1]


        # TODO Test that these are the correct images
        self.assertEquals(len(slide_with_image.images),1)

    def test_that_it_sets_page_offset(self):
        test_file_path = "src/test/test_data/offset-for-pages-test.pptx"

        coscrad_presentation = Presentation.from_pptx(test_file_path)

        third_slide = coscrad_presentation.slides[2]

        self.assertEquals(2002, third_slide.slide_number)

    def test_that_it_works_without_slide_offset(self):
        test_file_path = "src/test/test_data/test_presentation_with_images.pptx"

        coscrad_presentation = Presentation.from_pptx(test_file_path)

        third_slide = coscrad_presentation.slides[2]

        # Note that there is no offset here
        self.assertEquals(2, third_slide.slide_number)