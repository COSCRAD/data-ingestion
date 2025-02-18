import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from bs4 import BeautifulSoup
from zipfile import ZipFile


class PresentationSlide:
    def __init__(self, slide_id, slide_number):
        self.slide_id = slide_id

        self.slide_number = slide_number

        self.text_items = []

        self.images = []

    def add_image(self, pptx_image):
        self.images.append(pptx_image)

    def add_text(self, text):
        self.text_items.append(text)

    def shift_slide_number(self, offset):
        self.slide_number = self.slide_number + offset

    def __str__(self):
        return "\n".join(self.text_items)


def unzip_pptx(filepath):
    archive = ZipFile(filepath, "r")

    return archive


def parse_pptx_slide_number_offset(pptx_archive):
    # TODO verify that this file \ archive actually exists and fail fast if not
    xml_file = pptx_archive.open("ppt/presentation.xml")
    text = xml_file.read()
    bs = BeautifulSoup(text, "xml")

    el = bs.find("p:presentation")

    try:
        return int(el.attrs["firstSlideNum"])
    except:
        return 0


class Presentation:
    def __init__(self, slide_number_offset):
        self.slides = []

        # potentially None
        self.slide_number_offset = slide_number_offset

    def append_slide(self, slide):
        if self.slide_number_offset is not None:
            slide.shift_slide_number(self.slide_number_offset)

        self.slides.append(slide)

    def from_pptx(pptx_filepath):
        slide_offset = parse_pptx_slide_number_offset(unzip_pptx(pptx_filepath))

        pptx_presentation = pptx.Presentation(pptx_filepath)

        coscrad_presentation = Presentation(slide_number_offset=slide_offset)

        for slide in pptx_presentation.slides:
            coscrad_slide = PresentationSlide(
                slide_id=slide.slide_id, slide_number=len(coscrad_presentation.slides)
            )

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            coscrad_slide.add_text(run.text)

                # if(shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                # Note that `shape.shape_type` has been known to report a shape as text instead of an image incorrectly
                if hasattr(shape, "image"):
                    coscrad_slide.add_image(shape.image)

            coscrad_presentation.append_slide(coscrad_slide)

        return coscrad_presentation
